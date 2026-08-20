"""PPTX parser that extracts native text, tables, charts, images, and notes."""

import logging
from dataclasses import dataclass, field
from typing import List, Optional, Tuple
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

@dataclass
class ExtractedElement:
    element_type: str  # 'text', 'table', 'image', 'chart'
    top: int
    left: int
    content: Optional[str] = None  # Text or formatted markdown table
    image_bytes: Optional[bytes] = None

@dataclass
class SlideData:
    slide_index: int
    title: Optional[str] = None
    elements: List[ExtractedElement] = field(default_factory=list)
    speaker_notes: Optional[str] = None

def _extract_table_as_markdown(table) -> str:
    """Convert pptx table shape to Markdown table string."""
    rows = []
    num_cols = len(table.columns)
    if num_cols == 0:
        return ""

    for row in table.rows:
        row_vals = []
        for cell in row.cells:
            # Clean cell text: replace newlines with <br> or spaces
            val = cell.text.strip().replace("\n", "<br>").replace("|", "\\|")
            row_vals.append(val)
        rows.append(row_vals)

    if not rows:
        return ""

    lines = []
    # Header row
    lines.append("| " + " | ".join(rows[0]) + " |")
    # Separator
    lines.append("| " + " | ".join(["---"] * num_cols) + " |")
    # Data rows
    for r in rows[1:]:
        lines.append("| " + " | ".join(r) + " |")

    return "\n".join(lines)

def _extract_chart_text(chart) -> str:
    """Extract chart title, category names, series names."""
    parts = []
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            title = chart.chart_title.text_frame.text.strip()
            if title:
                parts.append(f"Chart: {title}")
    except Exception:
        pass

    try:
        series_names = [s.name for s in chart.series if getattr(s, 'name', None)]
        if series_names:
            parts.append(f"Series: {', '.join(series_names)}")
    except Exception:
        pass

    return "\n".join(parts)

def _process_shape(shape, elements_list: List[ExtractedElement], base_top: int = 0, base_left: int = 0):
    """Recursively inspect shapes, handling groups, textboxes, tables, images, charts."""
    shape_top = getattr(shape, "top", None)
    shape_left = getattr(shape, "left", None)
    top = (shape_top if shape_top is not None else 0) + (base_top if base_top is not None else 0)
    left = (shape_left if shape_left is not None else 0) + (base_left if base_left is not None else 0)

    # Handle Group Shapes recursively
    try:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP or hasattr(shape, "shapes"):
            for sub_shape in shape.shapes:
                _process_shape(sub_shape, elements_list, base_top=top, base_left=left)
            return
    except Exception as e:
        logger.debug(f"Group processing exception: {e}")

    # Handle Tables
    try:
        if getattr(shape, "has_table", False):
            md_table = _extract_table_as_markdown(shape.table)
            if md_table.strip():
                elements_list.append(ExtractedElement(
                    element_type="table",
                    top=top,
                    left=left,
                    content=md_table
                ))
            return
    except Exception as e:
        logger.warning(f"Error parsing table: {e}")

    # Handle Charts
    try:
        if getattr(shape, "has_chart", False):
            chart_text = _extract_chart_text(shape.chart)
            if chart_text.strip():
                elements_list.append(ExtractedElement(
                    element_type="chart",
                    top=top,
                    left=left,
                    content=chart_text
                ))
            return
    except Exception as e:
        logger.warning(f"Error parsing chart: {e}")

    # Handle Images / Pictures
    try:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
            image_bytes = shape.image.blob
            if image_bytes:
                elements_list.append(ExtractedElement(
                    element_type="image",
                    top=top,
                    left=left,
                    image_bytes=image_bytes
                ))
            return
    except Exception as e:
        logger.warning(f"Error extracting image: {e}")

    # Handle Text Frames
    try:
        if getattr(shape, "has_text_frame", False):
            raw_text = shape.text_frame.text.strip()
            # Filter out empty or placeholder prompts
            if raw_text and raw_text not in ["Click to edit Master title style", "Click to edit Master text styles"]:
                elements_list.append(ExtractedElement(
                    element_type="text",
                    top=top,
                    left=left,
                    content=raw_text
                ))
            return
    except Exception as e:
        logger.warning(f"Error extracting text frame: {e}")

def parse_pptx(file_path: str) -> Tuple[List[SlideData], int]:
    """
    Parse a .pptx file and return a list of SlideData objects and total slides count.
    """
    prs = pptx.Presentation(file_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides, start=1):
        elements: List[ExtractedElement] = []
        for shape in slide.shapes:
            _process_shape(shape, elements)

        # Sort elements by reading order: Top-to-Bottom, then Left-to-Right
        elements.sort(key=lambda elem: (elem.top, elem.left))

        # Extract Speaker Notes
        notes_text = None
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                raw_notes = slide.notes_slide.notes_text_frame.text.strip()
                if raw_notes:
                    notes_text = raw_notes
        except Exception as e:
            logger.debug(f"Slide {idx}: notes extraction skipped/failed: {e}")

        # Attempt to get title
        title_text = None
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text_frame.text.strip()
        except Exception:
            pass

        slide_data = SlideData(
            slide_index=idx,
            title=title_text,
            elements=elements,
            speaker_notes=notes_text
        )
        slides_data.append(slide_data)

    return slides_data, len(prs.slides)
